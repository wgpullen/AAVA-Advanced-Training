"""
AAVA Advanced Engineering Lab — PowerPoint Builder
Ascendion | COBOL Legacy Modernization Capstone | Employees Only
Build: June 2026

One-day (8-hour) instructor-led, MCP-driven course.
AM: core building blocks (decomposition, KB/GR/Tool, agent creation, governance).
PM: capstone — transform a real COBOL system (CUSTMGMT + ACCTPROC) to Spring Boot,
    retiring dead code with runtime-log evidence.

Output: AAVA_Advanced_Engineering_Lab_v1.pptx
Forked from build_aava_training.py (sales lab) — same Ascendion brand system.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Palette (Ascendion brand) ────────────────────────────────────────────────
ASC_DARK   = RGBColor(0x0C, 0x23, 0x40)
ASC_NAVY   = RGBColor(0x1A, 0x3A, 0x5C)
ASC_BLUE   = RGBColor(0x00, 0x78, 0xC8)
ASC_GREEN  = RGBColor(0x00, 0xA6, 0x51)
ASC_LIGHT  = RGBColor(0xE8, 0xF2, 0xFA)
ASC_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
ASC_GRAY   = RGBColor(0xF4, 0xF6, 0xFA)
ASC_TEXT   = RGBColor(0x1C, 0x1C, 0x2C)
ASC_MUTED  = RGBColor(0x6B, 0x72, 0x80)
ASC_BORDER = RGBColor(0xC8, 0xD8, 0xEC)
ASC_STEP   = RGBColor(0xE0, 0xF0, 0xFF)
ASC_RED    = RGBColor(0xC0, 0x39, 0x2B)
ASC_AMBER  = RGBColor(0xE6, 0x7E, 0x22)
ASC_SKY    = RGBColor(0xA0, 0xC0, 0xE0)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)
FOOTER  = "ASCENDION  ·  AAVA ADVANCED ENGINEERING  ·  EMPLOYEES ONLY"

# ── Core helpers (reused from the sales builder) ──────────────────────────────

def add_rect(slide, l, t, w, h, fill_rgb=None, line_rgb=None, line_width_pt=0):
    shape = slide.shapes.add_shape(1, l, t, w, h)
    shape.line.fill.background()
    if fill_rgb:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_rgb
    else:
        shape.fill.background()
    if line_rgb and line_width_pt > 0:
        shape.line.color.rgb = line_rgb
        shape.line.width = Pt(line_width_pt)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, l, t, w, h, size=12, bold=False, italic=False,
             color=ASC_TEXT, align=PP_ALIGN.LEFT, wrap=True, font="Calibri"):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font
    return box


def add_lines(slide, lines, l, t, w, h, size=11, bold=False, color=ASC_TEXT,
              align=PP_ALIGN.LEFT, spacing=None, font="Calibri"):
    """lines: list of str OR (text, bold, italic, size, color) tuples."""
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for item in lines:
        if isinstance(item, str):
            txt, b, it, sz, col = item, bold, False, size, color
        else:
            txt, b, it, sz, col = item
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        if spacing:
            p.line_spacing = Pt(spacing)
        if txt == "":
            p.add_run().text = ""
            continue
        run = p.add_run()
        run.text = txt
        run.font.size = Pt(sz)
        run.font.bold = b
        run.font.italic = it
        run.font.color.rgb = col
        run.font.name = font
    return box


def set_bg(slide, rgb):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def content_chrome(slide, module_label, title):
    set_bg(slide, ASC_WHITE)
    add_rect(slide, Inches(0), Inches(0), Inches(0.12), SLIDE_H, fill_rgb=ASC_NAVY)
    add_rect(slide, Inches(0.12), Inches(0), SLIDE_W - Inches(0.12), Inches(0.06), fill_rgb=ASC_BLUE)
    add_text(slide, module_label, Inches(0.25), Inches(0.08), Inches(9), Inches(0.28),
             size=7.5, bold=True, color=ASC_BLUE)
    add_text(slide, title, Inches(0.25), Inches(0.36), Inches(12.6), Inches(0.7),
             size=25, bold=True, color=ASC_DARK)
    add_rect(slide, Inches(0.25), Inches(1.08), Inches(12.83), Inches(0.03), fill_rgb=ASC_BORDER)
    add_text(slide, FOOTER, Inches(0.25), Inches(7.12), Inches(10), Inches(0.3),
             size=7, color=ASC_MUTED)


def section_header(prs, module_num, module_title, subtitle):
    slide = blank(prs)
    set_bg(slide, ASC_DARK)
    add_rect(slide, Inches(0), Inches(0), Inches(0.18), SLIDE_H, fill_rgb=ASC_GREEN)
    add_rect(slide, Inches(0.18), Inches(0), SLIDE_W - Inches(0.18), Inches(0.06), fill_rgb=ASC_BLUE)
    add_text(slide, module_num, Inches(0.5), Inches(1.6), Inches(11), Inches(0.55),
             size=13, bold=True, color=ASC_GREEN)
    add_text(slide, module_title, Inches(0.5), Inches(2.2), Inches(12), Inches(1.4),
             size=40, bold=True, color=ASC_WHITE)
    add_rect(slide, Inches(0.5), Inches(3.75), Inches(3), Inches(0.04), fill_rgb=ASC_BLUE)
    add_text(slide, subtitle, Inches(0.5), Inches(3.9), Inches(12), Inches(0.6),
             size=15, color=ASC_SKY)
    add_text(slide, FOOTER, Inches(0.5), Inches(7.12), Inches(10), Inches(0.3),
             size=7, color=RGBColor(0x50, 0x70, 0x90))
    return slide


def step_box(slide, num, text, x, y, w=Inches(5.9), tone=ASC_BLUE):
    box_h = Inches(0.52)
    add_rect(slide, x, y, w, box_h, fill_rgb=ASC_STEP, line_rgb=ASC_BORDER, line_width_pt=0.5)
    add_rect(slide, x, y, Inches(0.48), box_h, fill_rgb=tone)
    add_text(slide, str(num), x, y + Inches(0.06), Inches(0.48), Inches(0.38),
             size=14, bold=True, color=ASC_WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, text, x + Inches(0.56), y + Inches(0.06), w - Inches(0.64), Inches(0.42),
             size=10.5, color=ASC_TEXT)


def code_box(slide, lines_text, l, t, w, h, size=8.5):
    add_rect(slide, l, t, w, h, fill_rgb=RGBColor(0xF0, 0xF4, 0xF8),
             line_rgb=ASC_BORDER, line_width_pt=0.75)
    add_lines(slide, lines_text, l + Inches(0.15), t + Inches(0.1),
              w - Inches(0.3), h - Inches(0.2), size=size, color=ASC_NAVY, font="Courier New")


def card(slide, title, body_lines, x, y, w, h, accent=ASC_BLUE, title_color=None):
    add_rect(slide, x, y, w, h, fill_rgb=ASC_GRAY, line_rgb=ASC_BORDER, line_width_pt=0.5)
    add_rect(slide, x, y, w, Inches(0.06), fill_rgb=accent)
    add_text(slide, title, x + Inches(0.18), y + Inches(0.14), w - Inches(0.36), Inches(0.4),
             size=12.5, bold=True, color=title_color or ASC_DARK)
    add_lines(slide, body_lines, x + Inches(0.18), y + Inches(0.62),
              w - Inches(0.36), h - Inches(0.7), size=10, color=ASC_TEXT, spacing=14)


def callout(slide, text, x, y, w, tone=ASC_AMBER, h=Inches(0.75), label="⚠  GOTCHA"):
    add_rect(slide, x, y, w, h, fill_rgb=RGBColor(0xFD, 0xF3, 0xE6),
             line_rgb=tone, line_width_pt=1.0)
    add_rect(slide, x, y, Inches(0.08), h, fill_rgb=tone)
    add_text(slide, label, x + Inches(0.2), y + Inches(0.08), w - Inches(0.4), Inches(0.3),
             size=8.5, bold=True, color=tone)
    add_text(slide, text, x + Inches(0.2), y + Inches(0.36), w - Inches(0.4), h - Inches(0.42),
             size=9.5, color=ASC_TEXT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDES
# ══════════════════════════════════════════════════════════════════════════════

def s_cover(prs):
    slide = blank(prs)
    set_bg(slide, ASC_DARK)
    add_rect(slide, Inches(0), Inches(0), Inches(0.18), SLIDE_H, fill_rgb=ASC_GREEN)
    add_rect(slide, Inches(0.18), Inches(0), SLIDE_W - Inches(0.18), Inches(0.06), fill_rgb=ASC_BLUE)
    add_text(slide, "ASCENDION  ·  AAVA", Inches(0.5), Inches(0.22), Inches(10), Inches(0.4),
             size=10, bold=True, color=ASC_GREEN)
    add_lines(slide, [
        ("Advanced Engineering Lab", True, False, 46, ASC_WHITE),
        ("From Problem Statement to Production Multi-Agent Workflow", True, False, 19, ASC_BLUE),
    ], Inches(0.5), Inches(0.9), Inches(8.6), Inches(2.0))
    add_text(slide,
             "One day · 8 hours · Instructor-led · MCP-driven\n"
             "Capstone: COBOL → Spring Boot legacy modernization",
             Inches(0.5), Inches(3.0), Inches(8.0), Inches(1.0), size=14, color=ASC_SKY)
    add_rect(slide, Inches(0.5), Inches(4.05), Inches(3.5), Inches(0.03), fill_rgb=ASC_BLUE)
    add_lines(slide, [
        ("Ascendion Platform Engineering", True, False, 11, ASC_WHITE),
        ("June 2026  ·  Employees Only  ·  Built on AAVA-Mastery + AAVA-MCP-Server", False, False, 10, ASC_SKY),
    ], Inches(0.5), Inches(4.2), Inches(7.5), Inches(0.85))
    # Right panel — the day
    add_rect(slide, Inches(9.0), Inches(0.7), Inches(4.1), Inches(6.1),
             fill_rgb=RGBColor(0x15, 0x32, 0x52), line_rgb=ASC_BLUE, line_width_pt=0.75)
    add_text(slide, "YOUR DAY", Inches(9.15), Inches(0.95), Inches(3.8), Inches(0.4),
             size=9, bold=True, color=ASC_BLUE, align=PP_ALIGN.CENTER)
    rows = [
        ("MORNING", "Core building blocks", ASC_GREEN),
        ("M0–M2", "Setup · Commandments · Decomposition", ASC_WHITE),
        ("M3–M5", "KB/GR/Tool · Agents · Governance", ASC_WHITE),
        ("LUNCH", "60 minutes", ASC_AMBER),
        ("AFTERNOON", "Capstone — solve it", ASC_GREEN),
        ("M6–M10", "Brief · Design · Build · Test · Ship", ASC_WHITE),
    ]
    for i, (label, detail, col) in enumerate(rows):
        y = Inches(1.5) + i * Inches(0.84)
        add_rect(slide, Inches(9.15), y, Inches(3.78), Inches(0.68), fill_rgb=RGBColor(0x20, 0x40, 0x60))
        add_rect(slide, Inches(9.15), y, Inches(0.07), Inches(0.68), fill_rgb=col)
        add_text(slide, label, Inches(9.32), y + Inches(0.07), Inches(3.5), Inches(0.28),
                 size=9.5, bold=True, color=col)
        add_text(slide, detail, Inches(9.32), y + Inches(0.36), Inches(3.5), Inches(0.28),
                 size=8, color=ASC_SKY)
    add_text(slide, FOOTER, Inches(0.5), Inches(7.12), Inches(10), Inches(0.3),
             size=7, color=RGBColor(0x50, 0x70, 0x90))


def s_agenda(prs):
    slide = blank(prs)
    content_chrome(slide, "COURSE OVERVIEW", "The 8-Hour Day")
    am = [
        ("0:00", "M0 · Environment Bring-up", "Clone repos, setup.py, verify PAT + MCP"),
        ("0:30", "M1 · Mental Model + 11 Commandments", "The five primitives; the rules that prevent failure"),
        ("1:00", "M2 · Problem Decomposition", "One agent = one task; extract + decompose"),
        ("1:45", "M3 · KBs, Guardrails, Tools", "Build-vs-reuse; the 8-rule tool contract"),
        ("2:30", "M4 · Agent Creation via Repo + MCP", "Instruction template; push artifacts via MCP"),
        ("3:15", "M5 · Governance Discipline", "Test → Validate → Approve one-by-one"),
    ]
    pm = [
        ("4:30", "M6 · Capstone Briefing", "Meridian National Bank — the client problem"),
        ("4:50", "M7 · Design the Pipeline", "Apply Phase 1–2; lock your decomposition"),
        ("5:45", "M8 · Build via MCP", "KB → GR → Tool → Agent → Workflow"),
        ("7:00", "M9 · Test, Debug, Validate", "Unit-test agents; run the chain; debug"),
        ("7:40", "M10 · Approve & Debrief", "Ship it; defend your dead-code calls"),
    ]
    add_text(slide, "☀  MORNING — Core Building Blocks", Inches(0.4), Inches(1.25), Inches(6), Inches(0.35),
             size=13, bold=True, color=ASC_GREEN)
    for i, (t, title, sub) in enumerate(am):
        y = Inches(1.7) + i * Inches(0.78)
        add_rect(slide, Inches(0.4), y, Inches(6.1), Inches(0.68), fill_rgb=ASC_GRAY, line_rgb=ASC_BORDER, line_width_pt=0.5)
        add_text(slide, t, Inches(0.5), y + Inches(0.18), Inches(0.8), Inches(0.3), size=11, bold=True, color=ASC_BLUE)
        add_text(slide, title, Inches(1.3), y + Inches(0.07), Inches(5.1), Inches(0.3), size=10.5, bold=True, color=ASC_DARK)
        add_text(slide, sub, Inches(1.3), y + Inches(0.37), Inches(5.1), Inches(0.28), size=8.5, color=ASC_MUTED)
    add_text(slide, "🌆  AFTERNOON — Capstone", Inches(6.85), Inches(1.25), Inches(6), Inches(0.35),
             size=13, bold=True, color=ASC_GREEN)
    for i, (t, title, sub) in enumerate(pm):
        y = Inches(1.7) + i * Inches(0.78)
        add_rect(slide, Inches(6.85), y, Inches(6.1), Inches(0.68), fill_rgb=ASC_LIGHT, line_rgb=ASC_BORDER, line_width_pt=0.5)
        add_text(slide, t, Inches(6.95), y + Inches(0.18), Inches(0.8), Inches(0.3), size=11, bold=True, color=ASC_BLUE)
        add_text(slide, title, Inches(7.75), y + Inches(0.07), Inches(5.1), Inches(0.3), size=10.5, bold=True, color=ASC_DARK)
        add_text(slide, sub, Inches(7.75), y + Inches(0.37), Inches(5.1), Inches(0.28), size=8.5, color=ASC_MUTED)
    add_rect(slide, Inches(6.85), Inches(6.6), Inches(6.1), Inches(0.45), fill_rgb=ASC_DARK)
    add_text(slide, "🍽  Lunch 3:30–4:30 (60 min)", Inches(7.0), Inches(6.68), Inches(5.8), Inches(0.3),
             size=10, bold=True, color=ASC_WHITE)


def s_audience(prs):
    slide = blank(prs)
    content_chrome(slide, "BEFORE WE START", "Who This Is For & What You Need")
    card(slide, "WHO THIS IS FOR", [
        ("Ascendion delivery engineers, solution architects,", False, False, 10, ASC_TEXT),
        ("and delivery leads who will BUILD AAVA workflows", False, False, 10, ASC_TEXT),
        ("for clients — not the sales overview.", False, False, 10, ASC_TEXT),
        ("", False, False, 6, ASC_TEXT),
        ("You will leave able to take a raw client problem", False, False, 10, ASC_TEXT),
        ("and ship a production multi-agent workflow.", True, False, 10, ASC_NAVY),
    ], Inches(0.4), Inches(1.35), Inches(6.1), Inches(2.6), accent=ASC_GREEN)
    card(slide, "PREREQUISITES (set up before today)", [
        ("• Laptop: Python 3.9+, Git, Claude Code installed", False, False, 10, ASC_TEXT),
        ("• An AAVA PAT (bearer token) — realm 32 access", False, False, 10, ASC_TEXT),
        ("• Cloned: AAVA-Mastery, AAVA-MCP-Server", False, False, 10, ASC_TEXT),
        ("• Cloned: AAVA-Advanced-Training (this lab)", False, False, 10, ASC_TEXT),
        ("• MCP server 'mcp-aava' registered in Claude Code", False, False, 10, ASC_TEXT),
    ], Inches(6.85), Inches(1.35), Inches(6.1), Inches(2.6), accent=ASC_BLUE)
    add_rect(slide, Inches(0.4), Inches(4.25), Inches(12.55), Inches(2.5), fill_rgb=ASC_DARK)
    add_text(slide, "THE GOLDEN THREAD FOR THE WHOLE DAY", Inches(0.7), Inches(4.5), Inches(11.8), Inches(0.4),
             size=11, bold=True, color=ASC_GREEN)
    add_text(slide,
             "“AAVA is one agent per critical task. The hardest, most valuable skill is deciding\n"
             "what the tasks ARE — and proving your decisions with evidence.”",
             Inches(0.7), Inches(5.0), Inches(11.8), Inches(1.2), size=20, bold=True, color=ASC_WHITE)
    add_text(slide, "Everything this morning builds the skill the afternoon capstone tests.",
             Inches(0.7), Inches(6.25), Inches(11.8), Inches(0.4), size=12, italic=True, color=ASC_SKY)


# ── MODULE 0 ──────────────────────────────────────────────────────────────────
def s_m0_setup(prs):
    slide = blank(prs)
    content_chrome(slide, "MODULE 0 · ENVIRONMENT BRING-UP", "Get Green Before You Build")
    add_text(slide, "Everyone must reach a working connection before we teach a single concept.",
             Inches(0.4), Inches(1.25), Inches(12.5), Inches(0.4), size=12, italic=True, color=ASC_NAVY)
    steps = [
        "Clone AAVA-Mastery and AAVA-MCP-Server (GitHub or Azure DevOps — they stay in sync).",
        "Run  python setup.py  — it validates your PAT and resolves realm 32 / team 229.",
        "Confirm the MCP server 'mcp-aava' is registered in Claude Code (~/.claude.json).",
        "Run  test_aava_connection  → expect ✅  Realm: 32  ·  Token valid.",
        "Clone AAVA-Advanced-Training — the capstone lab artifacts live here.",
    ]
    for i, s in enumerate(steps):
        step_box(slide, i + 1, s, Inches(0.4), Inches(1.75) + i * Inches(0.62), w=Inches(7.4))
    code_box(slide, [
        "# verify you're wired up",
        "$ python setup.py",
        "$ # in Claude Code:",
        "  test_aava_connection()",
        "",
        "✅ AAVA CONNECTION SUCCESSFUL",
        "   User: <you>  Realm: 32",
        "   Status: Token valid",
    ], Inches(8.1), Inches(1.75), Inches(4.85), Inches(2.9))
    callout(slide,
            "Expired or wrong-realm PAT → setup.py re-prompts. Zscaler can block SSO, but the PAT "
            "path is what we use, so you're fine. Can't get MCP working? Pair up — but the course is "
            "designed for the MCP-driven path.",
            Inches(0.4), Inches(5.5), Inches(12.55), tone=ASC_AMBER, h=Inches(1.05))


def s_m0_surfaces(prs):
    slide = blank(prs)
    content_chrome(slide, "MODULE 0 · TOOLING", "Two AAVA Surfaces — Use Both")
    card(slide, "AAVA MCP Server  (primary today)", [
        ("63 RPC tools — flat API calls mid-conversation.", False, False, 10, ASC_TEXT),
        ("create_aava_* · execute_single_* · approve_aava_*", False, False, 9.5, ASC_NAVY),
        ("Fast, composable. This is how you'll author and", False, False, 10, ASC_TEXT),
        ("push every capstone artifact into AAVA.", False, False, 10, ASC_TEXT),
        ("", False, False, 6, ASC_TEXT),
        ("Tools appear as  mcp__mcp-aava__*  in Claude Code.", True, False, 9.5, ASC_BLUE),
    ], Inches(0.4), Inches(1.35), Inches(6.1), Inches(2.85), accent=ASC_BLUE)
    card(slide, "AAVA Claude Code Plugin  (/aava:*)", [
        ("Opinionated authoring + governance workflows.", False, False, 10, ASC_TEXT),
        ("/aava:author process · /aava:assess · /aava:approve", False, False, 9.5, ASC_NAVY),
        ("12 subagents, 14-dimension maturity assessor.", False, False, 10, ASC_TEXT),
        ("", False, False, 6, ASC_TEXT),
        ("Use when you want the maturity model end-to-end.", True, False, 9.5, ASC_GREEN),
    ], Inches(6.85), Inches(1.35), Inches(6.1), Inches(2.85), accent=ASC_GREEN)
    add_rect(slide, Inches(0.4), Inches(4.5), Inches(12.55), Inches(2.2), fill_rgb=ASC_LIGHT, line_rgb=ASC_BORDER, line_width_pt=0.5)
    add_text(slide, "DECISION RULE", Inches(0.7), Inches(4.7), Inches(11), Inches(0.35), size=11, bold=True, color=ASC_DARK)
    add_lines(slide, [
        ("• Need a specific API call right now  →  MCP tool   (e.g. list_aava_artifacts, create_aava_agent)", False, False, 11, ASC_TEXT),
        ("• Need to design a pipeline / govern through lifecycle / score the 14-dim rubric  →  Plugin skill", False, False, 11, ASC_TEXT),
        ("• Both surfaces hit the same AAVA platform with the same PAT.", True, False, 11, ASC_NAVY),
    ], Inches(0.7), Inches(5.1), Inches(11.9), Inches(1.4), spacing=20)


# ── MODULE 1 ──────────────────────────────────────────────────────────────────
def s_m1_glance(prs):
    slide = blank(prs)
    content_chrome(slide, "MODULE 1 · MENTAL MODEL", "AAVA at a Glance — Five Primitives")
    prims = [
        ("Knowledge Base", "Documents agents RAG against", ASC_BLUE),
        ("Guardrail", "Constraints on agent output", ASC_AMBER),
        ("Tool", "A Python action the LLM can call", ASC_GREEN),
        ("Agent", "An LLM with role, goal, tools, KB", ASC_NAVY),
        ("Workflow", "Sequential pipeline of agents", ASC_BLUE),
    ]
    for i, (name, desc, col) in enumerate(prims):
        x = Inches(0.4) + i * Inches(2.55)
        add_rect(slide, x, Inches(1.45), Inches(2.4), Inches(1.7), fill_rgb=ASC_GRAY, line_rgb=col, line_width_pt=1.0)
        add_rect(slide, x, Inches(1.45), Inches(2.4), Inches(0.45), fill_rgb=col)
        add_text(slide, name, x, Inches(1.5), Inches(2.4), Inches(0.4), size=12, bold=True, color=ASC_WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, desc, x + Inches(0.12), Inches(2.05), Inches(2.16), Inches(1.0), size=10, color=ASC_TEXT, align=PP_ALIGN.CENTER)
        if i < 4:
            add_text(slide, "→", x + Inches(2.42), Inches(2.05), Inches(0.2), Inches(0.4), size=16, bold=True, color=ASC_MUTED)
    add_rect(slide, Inches(0.4), Inches(3.5), Inches(12.55), Inches(1.2), fill_rgb=ASC_LIGHT, line_rgb=ASC_BORDER, line_width_pt=0.5)
    add_text(slide, "REALM = the org boundary (team / project scope)", Inches(0.7), Inches(3.65), Inches(12), Inches(0.35),
             size=12, bold=True, color=ASC_DARK)
    add_text(slide,
             "You build in realm 32 (platformengineeringallteam), team 229. Realms 1 (Ascendion root), "
             "59 (Executive), 75 (asc-markets-all) also exist. Artifacts are scoped to a realm.",
             Inches(0.7), Inches(4.05), Inches(11.9), Inches(0.6), size=10.5, color=ASC_TEXT)
    add_text(slide, "Built on the CrewAI framework. Platform: https://int-ai.aava.ai",
             Inches(0.4), Inches(5.0), Inches(12), Inches(0.4), size=11, italic=True, color=ASC_NAVY)
    callout(slide,
            "An Agent is only as good as its decomposition. Most failures aren't model failures — "
            "they're scoping failures. Keep that in mind through Module 2.",
            Inches(0.4), Inches(5.6), Inches(12.55), tone=ASC_BLUE, h=Inches(0.95), label="◆  KEY IDEA")


def s_m1_commandments_a(prs):
    slide = blank(prs)
    content_chrome(slide, "MODULE 1 · THE 11 COMMANDMENTS", "Scar Tissue From 16 Client Projects (1–6)")
    cmds = [
        ("1", "NEVER put guardrails on workflows — only on agents.", ASC_RED),
        ("2", "azure_parent_folder MUST propagate through ALL agents.", ASC_RED),
        ("3", "Tool names must match EXACTLY (spaces, 'Tool' suffix).", ASC_AMBER),
        ("4", "AzureBlobWriterTool param is 'content', not file_content.", ASC_AMBER),
        ("5", "Upload code files BEFORE generating the HTML report.", ASC_BLUE),
        ("6", "HTML reports use the ACTUAL SAS URLs — never construct them.", ASC_BLUE),
    ]
    for i, (n, t, col) in enumerate(cmds):
        y = Inches(1.45) + i * Inches(0.86)
        add_rect(slide, Inches(0.4), y, Inches(12.55), Inches(0.74), fill_rgb=ASC_GRAY, line_rgb=ASC_BORDER, line_width_pt=0.5)
        add_rect(slide, Inches(0.4), y, Inches(0.7), Inches(0.74), fill_rgb=col)
        add_text(slide, n, Inches(0.4), y + Inches(0.16), Inches(0.7), Inches(0.45), size=20, bold=True, color=ASC_WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, t, Inches(1.3), y + Inches(0.2), Inches(11.4), Inches(0.4), size=13, bold=True, color=ASC_DARK)
    add_text(slide, "Red = breaks the chain silently · Amber = silent wrong output · Blue = report integrity",
             Inches(0.4), Inches(6.75), Inches(12.5), Inches(0.3), size=8.5, italic=True, color=ASC_MUTED)


def s_m1_commandments_b(prs):
    slide = blank(prs)
    content_chrome(slide, "MODULE 1 · THE 11 COMMANDMENTS", "Scar Tissue From 16 Client Projects (7–11)")
    cmds = [
        ("7", "Agents fake data instead of calling tools (the silent failure).", ASC_RED,
         "Fix: merge 'generate content' + 'call tool' into ONE mandatory step."),
        ("8", "Use 'DO NOT PROCEED' blocking language on every tool call.", ASC_RED,
         "Without it, agents skip tool calls 60%+ of the time."),
        ("9", "No 'Section' numbering inside Steps.", ASC_AMBER,
         "Steps 1..N for actions; PART A..G for sub-items. Mixing confuses the agent."),
        ("10", "Creation order is ALWAYS KB → GR → Tool → Agent → Workflow.", ASC_GREEN,
         "Each layer references the one before it. Build bottom-up."),
        ("11", "Clone an artifact ONLY when you must modify one you don't own / is APPROVED.", ASC_BLUE,
         "APPROVED is immutable; clone via POST /agents/clone, then edit the clone."),
    ]
    for i, (n, t, col, sub) in enumerate(cmds):
        y = Inches(1.45) + i * Inches(1.04)
        add_rect(slide, Inches(0.4), y, Inches(12.55), Inches(0.92), fill_rgb=ASC_GRAY, line_rgb=ASC_BORDER, line_width_pt=0.5)
        add_rect(slide, Inches(0.4), y, Inches(0.7), Inches(0.92), fill_rgb=col)
        add_text(slide, n, Inches(0.4), y + Inches(0.26), Inches(0.7), Inches(0.45), size=18, bold=True, color=ASC_WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, t, Inches(1.3), y + Inches(0.12), Inches(11.4), Inches(0.4), size=12.5, bold=True, color=ASC_DARK)
        add_text(slide, sub, Inches(1.3), y + Inches(0.5), Inches(11.4), Inches(0.36), size=10, italic=True, color=ASC_NAVY)


# ── MODULE 2 ──────────────────────────────────────────────────────────────────
def s_m2_golden(prs):
    slide = blank(prs)
    content_chrome(slide, "MODULE 2 · DECOMPOSITION", "One Agent = One Independent Task")
    add_rect(slide, Inches(0.4), Inches(1.35), Inches(12.55), Inches(1.3), fill_rgb=ASC_DARK)
    add_text(slide, "THE GOLDEN RULE", Inches(0.7), Inches(1.5), Inches(11), Inches(0.35), size=11, bold=True, color=ASC_GREEN)
    add_text(slide, "If an agent does two things, split it. If you can't test it alone, it's too big.",
             Inches(0.7), Inches(1.9), Inches(11.9), Inches(0.6), size=19, bold=True, color=ASC_WHITE)
    card(slide, "✅  Good — testable, debuggable, reusable", [
        ("Agent 1: parse the COBOL (structure only)", False, False, 10.5, ASC_TEXT),
        ("Agent 2: analyze runtime logs (active vs dead)", False, False, 10.5, ASC_TEXT),
        ("Agent 3: document the active code", False, False, 10.5, ASC_TEXT),
        ("Each agent = one thing that can fail =", True, False, 10.5, ASC_GREEN),
        ("one thing you can fix in isolation.", True, False, 10.5, ASC_GREEN),
    ], Inches(0.4), Inches(2.85), Inches(6.1), Inches(2.4), accent=ASC_GREEN)
    card(slide, "❌  Rookie move — the mega-agent", [
        ("Agent 1: 'analyze the COBOL, find dead code,", False, False, 10.5, ASC_TEXT),
        ("document it, write stories, design the target,", False, False, 10.5, ASC_TEXT),
        ("generate Java, write tests, and report.'", False, False, 10.5, ASC_TEXT),
        ("Can't test. Can't debug. Can't reuse.", True, False, 10.5, ASC_RED),
        ("Fails opaquely halfway through.", True, False, 10.5, ASC_RED),
    ], Inches(6.85), Inches(2.85), Inches(6.1), Inches(2.4), accent=ASC_RED)
    callout(slide,
            "Automatability classification helps size the work: >80% / 50–80% / 20–50% / <20%. "
            "But the first cut is always: what are the INDEPENDENT tasks?",
            Inches(0.4), Inches(5.5), Inches(12.55), tone=ASC_BLUE, h=Inches(0.9), label="◆  HEURISTIC")


def s_m2_phase1(prs):
    slide = blank(prs)
    content_chrome(slide, "MODULE 2 · PHASE 1", "Extract the 8 Critical Elements")
    add_text(slide, "From any raw client input (email, RFP, meeting notes), pull these out before you design anything:",
             Inches(0.4), Inches(1.25), Inches(12.5), Inches(0.4), size=12, italic=True, color=ASC_NAVY)
    elems = [
        ("1  Business goal", "What outcome does the client actually want?"),
        ("2  Inputs provided", "What data do you have to work with?"),
        ("3  Outputs required", "What concrete deliverables must you produce?"),
        ("4  Constraints", "What limits the solution (cost, policy, scope)?"),
        ("5  Domain rules", "What business logic / terminology applies?"),
        ("6  Risk / compliance", "What can't go wrong? Who's watching?"),
        ("7  Success metric", "How will the client measure success?"),
        ("8  Hidden traps", "What looks obvious but isn't? (find these!)"),
    ]
    for i, (h, d) in enumerate(elems):
        col = i % 2
        row = i // 2
        x = Inches(0.4) + col * Inches(6.45)
        y = Inches(1.85) + row * Inches(1.15)
        accent = ASC_RED if i == 7 else ASC_BLUE
        add_rect(slide, x, y, Inches(6.2), Inches(1.0), fill_rgb=ASC_GRAY, line_rgb=ASC_BORDER, line_width_pt=0.5)
        add_rect(slide, x, y, Inches(0.08), Inches(1.0), fill_rgb=accent)
        add_text(slide, h, x + Inches(0.22), y + Inches(0.12), Inches(5.9), Inches(0.35), size=13, bold=True, color=ASC_DARK)
        add_text(slide, d, x + Inches(0.22), y + Inches(0.5), Inches(5.9), Inches(0.4), size=10.5, color=ASC_TEXT)
    add_text(slide, "Element 8 is where the value is. This afternoon's capstone has a deliberate hidden trap.",
             Inches(0.4), Inches(6.75), Inches(12.5), Inches(0.3), size=9.5, bold=True, italic=True, color=ASC_RED)


def s_m2_phase2(prs):
    slide = blank(prs)
    content_chrome(slide, "MODULE 2 · PHASE 2", "Decompose Into a Pipeline")
    card(slide, "PROVEN PIPELINE SHAPE (Golden Pattern)", [
        ("Generator  →  QA / Analyzer  →  Finalizer  →  Output Creator", True, False, 11, ASC_NAVY),
        ("", False, False, 4, ASC_TEXT),
        ("• First agent reads input + generates azure_parent_folder", False, False, 10, ASC_TEXT),
        ("• Middle agents each do ONE transformation, read prev / write own", False, False, 10, ASC_TEXT),
        ("• Final agent produces the customer-facing artifact (exec report)", False, False, 10, ASC_TEXT),
        ("• HITL gates placed by risk tier between stages where needed", False, False, 10, ASC_TEXT),
    ], Inches(0.4), Inches(1.35), Inches(7.6), Inches(2.7), accent=ASC_GREEN)
    card(slide, "HOW MANY AGENTS?", [
        ("Count the independent", False, False, 10.5, ASC_TEXT),
        ("transformations between", False, False, 10.5, ASC_TEXT),
        ("raw input and final output.", False, False, 10.5, ASC_TEXT),
        ("", False, False, 4, ASC_TEXT),
        ("Each becomes one agent.", True, False, 10.5, ASC_BLUE),
        ("Typical: 3–8 agents.", True, False, 10.5, ASC_BLUE),
    ], Inches(8.15), Inches(1.35), Inches(4.8), Inches(2.7), accent=ASC_BLUE)
    add_rect(slide, Inches(0.4), Inches(4.35), Inches(12.55), Inches(2.35), fill_rgb=ASC_LIGHT, line_rgb=ASC_BORDER, line_width_pt=0.5)
    add_text(slide, "🛠  DRILL (15 min): decompose 'Generate release notes from merged PRs'",
             Inches(0.7), Inches(4.55), Inches(12), Inches(0.4), size=12, bold=True, color=ASC_DARK)
    add_lines(slide, [
        ("Whiteboard your agent breakdown, then compare. A clean answer:", False, False, 11, ASC_TEXT),
        ("   Agent 1  PR / commit extractor   →   Agent 2  release-notes writer   →   Agent 3  formatter / publisher", True, False, 12, ASC_NAVY),
        ("Three independently testable tasks. Notice each could be unit-tested with execute_single_agent.", False, True, 10.5, ASC_MUTED),
    ], Inches(0.7), Inches(5.05), Inches(11.9), Inches(1.5), spacing=22)


# ── MODULE 3 ──────────────────────────────────────────────────────────────────
def s_m3_kb_gr(prs):
    slide = blank(prs)
    content_chrome(slide, "MODULE 3 · KNOWLEDGE & RULES", "Knowledge Bases and Guardrails")
    card(slide, "KNOWLEDGE BASE — shared source of truth", [
        ("Put in a KB (not the prompt): regulatory refs,", False, False, 10, ASC_TEXT),
        ("schemas, mappings, long lists, client domain rules.", False, False, 10, ASC_TEXT),
        ("", False, False, 4, ASC_TEXT),
        ("Model: cobol-modernization-standards encodes the", False, False, 10, ASC_TEXT),
        ("COBOL→Spring mapping + dead-code criteria so all", False, False, 10, ASC_TEXT),
        ("8 agents share ONE definition of 'dead code'.", True, False, 10, ASC_NAVY),
        ("", False, False, 4, ASC_TEXT),
        ("Create:  create_aava_knowledge_base", True, False, 9.5, ASC_BLUE),
    ], Inches(0.4), Inches(1.35), Inches(6.1), Inches(3.5), accent=ASC_BLUE)
    card(slide, "GUARDRAIL — rules agents must always follow", [
        ("Tier-1 (platform-managed): PII, Compliance,", False, False, 10, ASC_TEXT),
        ("Toxicity, Enterprise Compliance — don't author these.", False, False, 10, ASC_TEXT),
        ("Tier-2 (domain): you author these.", False, False, 10, ASC_TEXT),
        ("", False, False, 4, ASC_TEXT),
        ("Colang = exact keyword match.", False, False, 10, ASC_TEXT),
        ("YAML  = semantic / fuzzy detection.", False, False, 10, ASC_TEXT),
        ("", False, False, 4, ASC_TEXT),
        ("⛔ Commandment #1: attach to AGENTS, never workflows.", True, False, 9.5, ASC_RED),
    ], Inches(6.85), Inches(1.35), Inches(6.1), Inches(3.5), accent=ASC_AMBER)
    callout(slide,
            "Micro-lab: create a tiny KB via create_aava_knowledge_base, then find_aava_artifact the "
            "AzureBlob tools so the whole room sees reuse-before-create before we touch agents.",
            Inches(0.4), Inches(5.15), Inches(12.55), tone=ASC_GREEN, h=Inches(1.0), label="🛠  MICRO-LAB")


def s_m3_tools(prs):
    slide = blank(prs)
    content_chrome(slide, "MODULE 3 · TOOLS", "The 8-Rule Tool Contract + Reuse")
    rules = [
        "No secrets in source — use K8s env / Config Service.",
        "stdlib only (no third-party imports in tool source).",
        "Schema = LLM-visible inputs only.",
        "Right HTTP shape per endpoint.",
        "Dual-status check on poll.",
        "Hyphens in variable names.",
        "Full toolConfig on PUT.",
        "Single responsibility — one tool, one job.",
    ]
    for i, r in enumerate(rules):
        col = i % 2
        row = i // 2
        x = Inches(0.4) + col * Inches(6.45)
        y = Inches(1.45) + row * Inches(0.72)
        add_rect(slide, x, y, Inches(6.2), Inches(0.6), fill_rgb=ASC_GRAY, line_rgb=ASC_BORDER, line_width_pt=0.5)
        add_text(slide, f"{i+1}", x + Inches(0.1), y + Inches(0.13), Inches(0.5), Inches(0.35), size=14, bold=True, color=ASC_GREEN)
        add_text(slide, r, x + Inches(0.65), y + Inches(0.15), Inches(5.4), Inches(0.35), size=10.5, color=ASC_TEXT)
    add_rect(slide, Inches(0.4), Inches(4.5), Inches(12.55), Inches(2.2), fill_rgb=ASC_DARK)
    add_text(slide, "REUSE > CREATE — the shared Azure tools already exist in realm 32",
             Inches(0.7), Inches(4.68), Inches(12), Inches(0.4), size=12, bold=True, color=ASC_GREEN)
    add_lines(slide, [
        ("Current DateTime Tool (4525)      generate azure_parent_folder", False, False, 11, ASC_WHITE),
        ("AzureBlobReaderTool (4521)        read one file", False, False, 11, ASC_WHITE),
        ("AzureBlobWriterTool (5964)        write file + return SAS URL", False, False, 11, ASC_WHITE),
        ("AzureBlobRecursiveReader (3412)   read all files in a folder", False, False, 11, ASC_WHITE),
        ("Always find_aava_artifact FIRST. Never recreate a tool that exists.", True, False, 11, ASC_SKY),
    ], Inches(0.7), Inches(5.1), Inches(11.9), Inches(1.5), spacing=18, font="Courier New")


# ── MODULE 4 ──────────────────────────────────────────────────────────────────
def s_m4_template(prs):
    slide = blank(prs)
    content_chrome(slide, "MODULE 4 · AGENT CREATION", "The Proven Instruction Template")
    code_box(slide, [
        "You are Agent #N in an M-agent pipeline.",
        "You receive azure_parent_folder from Agent #N-1.",
        "",
        "STEP 1 — <action> (MANDATORY TOOL CALL)",
        "  YOU MUST call \"<Exact Tool Name>\" with: <params>",
        "  DO NOT PROCEED until the tool returns <result>.",
        "",
        "STEP K — Upload via \"AzureBlobWriterTool\"",
        "  content = <your markdown>",
        "  DO NOT PROCEED until SAS URL returned.",
        "",
        "STEP K+1 — Output handoff:",
        "  azure_parent_folder + SAS URL + summary for #N+1",
    ], Inches(0.4), Inches(1.35), Inches(7.3), Inches(4.0), size=10)
    card(slide, "PROMPT STRUCTURE", [
        ("Role → Goal → Context →", True, False, 11, ASC_NAVY),
        ("Output → Rules → Example", True, False, 11, ASC_NAVY),
        ("", False, False, 4, ASC_TEXT),
        ("• 150–300 words is the sweet spot", False, False, 10, ASC_TEXT),
        ("• Positive instructions (what TO do)", False, False, 10, ASC_TEXT),
        ("• Blocking language on tool calls", False, False, 10, ASC_TEXT),
        ("• Description field ≤ ~100 lines", False, False, 10, ASC_TEXT),
    ], Inches(7.9), Inches(1.35), Inches(5.05), Inches(2.55), accent=ASC_BLUE)
    callout(slide,
            "See Agents 1 & 2 in solution-key/reference-agents/. Agent 1 (static parse) and Agent 2 "
            "(runtime log analysis) are the model for the STEP / DO-NOT-PROCEED pattern.",
            Inches(7.9), Inches(4.05), Inches(5.05), tone=ASC_GREEN, h=Inches(1.3), label="📄  REFERENCE")
    callout(slide,
            "Commandment #7 + #8 live here. Agents 'helpfully' summarize data they never read. "
            "Merge generate+call into one step and gate it with DO NOT PROCEED.",
            Inches(0.4), Inches(5.55), Inches(7.3), tone=ASC_RED, h=Inches(1.0))


def s_m4_configs_mcp(prs):
    slide = blank(prs)
    content_chrome(slide, "MODULE 4 · CONFIG + PUSH VIA MCP", "agentConfigs & create_aava_agent")
    code_box(slide, [
        "\"agentConfigs\": {",
        "  \"aiEngine\": \"AzureOpenAI\",",
        "  \"model\": \"gpt-4o\", \"modelId\": 53,",
        "  \"preset\": \"Deterministic\",",
        "  \"temperature\": 0.1, \"topP\": 0.6,",
        "  \"maxIter\": 3, \"maxRpm\": 10,",
        "  \"maxExecutionTime\": 1500",
        "}",
    ], Inches(0.4), Inches(1.4), Inches(6.3), Inches(2.7), size=10.5)
    card(slide, "MODEL SELECTION", [
        ("Technical / parse / code:", True, False, 10.5, ASC_NAVY),
        ("  gpt-4o (53) or Bedrock Claude (497/493)", False, False, 10, ASC_TEXT),
        ("Reports / creative: Balanced preset", False, False, 10, ASC_TEXT),
        ("Deterministic = code/classification (temp 0.1)", False, False, 10, ASC_TEXT),
        ("", False, False, 3, ASC_TEXT),
        ("Reference pipeline uses gpt-4o throughout.", True, False, 10, ASC_BLUE),
    ], Inches(6.95), Inches(1.4), Inches(6.0), Inches(2.7), accent=ASC_BLUE)
    add_rect(slide, Inches(0.4), Inches(4.35), Inches(12.55), Inches(2.35), fill_rgb=ASC_DARK)
    add_text(slide, "PUSH TO AAVA VIA MCP", Inches(0.7), Inches(4.55), Inches(11), Inches(0.4), size=12, bold=True, color=ASC_GREEN)
    add_lines(slide, [
        ("create_aava_agent(...)   →  lands as CREATED/DRAFT in realm 32", False, False, 11, ASC_WHITE),
        ("LLM params nest inside agentConfigs (the Jun-2026 MCP fix — top-level = HTTP 500).", False, False, 11, ASC_SKY),
        ("🛡 Revelio Safeguard: create_aava_agent auto-restores description + expectedOutput", False, False, 11, ASC_WHITE),
        ("   if Revelio overwrites them post-creation. Silent — you don't lose your text.", False, False, 11, ASC_WHITE),
    ], Inches(0.7), Inches(5.0), Inches(11.9), Inches(1.5), spacing=20, font="Courier New")


# ── MODULE 5 ──────────────────────────────────────────────────────────────────
def s_m5_governance(prs):
    slide = blank(prs)
    content_chrome(slide, "MODULE 5 · GOVERNANCE DISCIPLINE", "Test → Validate → Approve One-by-One")
    flow = [
        ("BUILD", "create_aava_* pushes the artifact", ASC_BLUE),
        ("TEST", "execute_single_agent / execute_single_tool", ASC_AMBER),
        ("VALIDATE", "read the output — is it actually correct?", ASC_GREEN),
        ("APPROVE", "approve_aava_artifact — one at a time", ASC_NAVY),
    ]
    for i, (label, desc, col) in enumerate(flow):
        x = Inches(0.4) + i * Inches(3.2)
        add_rect(slide, x, Inches(1.55), Inches(3.0), Inches(1.7), fill_rgb=ASC_GRAY, line_rgb=col, line_width_pt=1.0)
        add_rect(slide, x, Inches(1.55), Inches(3.0), Inches(0.5), fill_rgb=col)
        add_text(slide, label, x, Inches(1.62), Inches(3.0), Inches(0.4), size=14, bold=True, color=ASC_WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, desc, x + Inches(0.15), Inches(2.2), Inches(2.7), Inches(0.95), size=9.5, color=ASC_TEXT, align=PP_ALIGN.CENTER)
        if i < 3:
            add_text(slide, "→", x + Inches(3.02), Inches(2.2), Inches(0.18), Inches(0.4), size=18, bold=True, color=ASC_MUTED)
    add_rect(slide, Inches(0.4), Inches(3.7), Inches(12.55), Inches(1.6), fill_rgb=ASC_LIGHT, line_rgb=ASC_BORDER, line_width_pt=0.5)
    add_text(slide, "THE ONE RULE", Inches(0.7), Inches(3.88), Inches(11), Inches(0.35), size=11, bold=True, color=ASC_RED)
    add_text(slide, "Never approve an artifact you haven't seen produce correct output.",
             Inches(0.7), Inches(4.28), Inches(11.9), Inches(0.5), size=18, bold=True, color=ASC_DARK)
    add_text(slide, "Super Admins can self-approve. Approval body uses  id  (not agentId). Approve in creation order.",
             Inches(0.7), Inches(4.85), Inches(11.9), Inches(0.4), size=11, color=ASC_NAVY)
    callout(slide,
            "This is the exact loop you'll run in M9–M10 on the capstone. Build all the artifacts, "
            "test each one, validate the output, then approve them one-by-one. Set the habit now.",
            Inches(0.4), Inches(5.55), Inches(12.55), tone=ASC_BLUE, h=Inches(1.0), label="◆  BRIDGE TO THE CAPSTONE")


def s_lunch(prs):
    slide = blank(prs)
    set_bg(slide, ASC_DARK)
    add_rect(slide, Inches(0), Inches(0), Inches(0.18), SLIDE_H, fill_rgb=ASC_GREEN)
    add_text(slide, "🍽", Inches(0.5), Inches(2.0), Inches(3), Inches(1.5), size=72, color=ASC_WHITE)
    add_text(slide, "Lunch — 60 minutes", Inches(0.5), Inches(3.5), Inches(12), Inches(1.0),
             size=44, bold=True, color=ASC_WHITE)
    add_text(slide, "Back at the top of the hour. This afternoon you SOLVE a real client problem.",
             Inches(0.5), Inches(4.6), Inches(12), Inches(0.5), size=16, color=ASC_SKY)
    add_text(slide, FOOTER, Inches(0.5), Inches(7.12), Inches(10), Inches(0.3), size=7, color=RGBColor(0x50, 0x70, 0x90))


# ── MODULE 6 — CAPSTONE BRIEFING ──────────────────────────────────────────────
def s_m6_client(prs):
    slide = blank(prs)
    content_chrome(slide, "MODULE 6 · CAPSTONE BRIEFING", "Client: Meridian National Bank")
    add_rect(slide, Inches(0.4), Inches(1.3), Inches(12.55), Inches(2.85), fill_rgb=RGBColor(0xFB, 0xFC, 0xFE), line_rgb=ASC_BORDER, line_width_pt=0.75)
    add_text(slide, "From: Eleanor Vance, SVP Core Platforms   ·   Subject: Mainframe modernization (board ask)",
             Inches(0.6), Inches(1.45), Inches(12.2), Inches(0.4), size=10.5, bold=True, color=ASC_BLUE)
    add_lines(slide, [
        ("“The board approved funding to get us OFF the mainframe for our Customer & Account core.", False, True, 12, ASC_TEXT),
        ("Two COBOL programs run this: CUSTMGMT (customer CRUD + reporting) and ACCTPROC", False, True, 12, ASC_TEXT),
        ("(nightly account batch). We want it on Spring Boot / Java.", False, True, 12, ASC_TEXT),
        ("", False, False, 4, ASC_TEXT),
        ("The previous vendor quoted a 1-for-1 rewrite of EVERY line at $2.4M / 14 months. Half this", False, True, 12, ASC_TEXT),
        ("code looks like it hasn't run in years — but nobody will sign off on deleting anything", False, True, 12, ASC_TEXT),
        ("because we can't PROVE what's safe to drop. Don't make me pay to rebuild dead code —", False, True, 12, ASC_TEXT),
        ("but don't delete anything you can't justify with data. Show me your reasoning.”", False, True, 12, ASC_NAVY),
    ], Inches(0.6), Inches(1.9), Inches(12.2), Inches(2.1), spacing=18)
    card(slide, "WHAT THE BOARD WANTS (in order)", [
        ("A.  A defensible plan: migrate / RETIRE + evidence", False, False, 10.5, ASC_TEXT),
        ("B.  Modern docs + user stories for what we keep", False, False, 10.5, ASC_TEXT),
        ("C.  Spring Boot design + code + tests (active logic)", False, False, 10.5, ASC_TEXT),
        ("D.  Executive summary for the board", False, False, 10.5, ASC_TEXT),
    ], Inches(0.4), Inches(4.35), Inches(7.6), Inches(2.4), accent=ASC_GREEN)
    card(slide, "THE CONSTRAINT", [
        ("Don't rebuild", True, False, 12, ASC_NAVY),
        ("untouched code.", True, False, 12, ASC_NAVY),
        ("", False, False, 3, ASC_TEXT),
        ("But prove EVERY", True, False, 12, ASC_RED),
        ("deletion with data.", True, False, 12, ASC_RED),
    ], Inches(8.15), Inches(4.35), Inches(4.8), Inches(2.4), accent=ASC_RED)


def s_m6_artifacts(prs):
    slide = blank(prs)
    content_chrome(slide, "MODULE 6 · CAPSTONE", "What the Client Handed You")
    arts = [
        ("source/CUSTMGMT.cbl", "Customer mgmt — online (CICS) + batch", ASC_BLUE),
        ("source/ACCTPROC.cbl", "Nightly account-processing batch", ASC_BLUE),
        ("source/copybooks/CUSTREC.cpy", "Customer record layout", ASC_GREEN),
        ("source/copybooks/ACCTREC.cpy", "Account record layout", ASC_GREEN),
        ("logs/batch_jcl_execution.log", "Full-year 2024 JCL batch log + exec summary", ASC_AMBER),
        ("logs/cics_transaction.log", "Full-year 2024 CICS transaction log", ASC_AMBER),
    ]
    for i, (path, desc, col) in enumerate(arts):
        y = Inches(1.4) + i * Inches(0.66)
        add_rect(slide, Inches(0.4), y, Inches(12.55), Inches(0.56), fill_rgb=ASC_GRAY, line_rgb=ASC_BORDER, line_width_pt=0.5)
        add_rect(slide, Inches(0.4), y, Inches(0.08), Inches(0.56), fill_rgb=col)
        add_text(slide, path, Inches(0.6), y + Inches(0.13), Inches(5.0), Inches(0.32), size=11, bold=True, color=ASC_DARK, font="Courier New")
        add_text(slide, desc, Inches(5.7), y + Inches(0.14), Inches(7.0), Inches(0.32), size=10.5, color=ASC_TEXT)
    add_rect(slide, Inches(0.4), Inches(5.5), Inches(12.55), Inches(1.2), fill_rgb=ASC_DARK)
    add_text(slide, "TWO WAYS TO FEED YOUR AGENTS", Inches(0.7), Inches(5.65), Inches(12), Inches(0.35), size=11, bold=True, color=ASC_GREEN)
    add_lines(slide, [
        ("① Azure Blob (primary):  container  cobol-legacy-modernization  →  folders  source/  and  logs/", False, False, 11, ASC_WHITE),
        ("② GitHub raw (fallback):  raw.githubusercontent.com/wgpullen/AAVA-Advanced-Training/main/lab-artifacts/", False, False, 11, ASC_SKY),
    ], Inches(0.7), Inches(6.05), Inches(12.2), Inches(0.6), spacing=18, font="Courier New")


# ── MODULE 7 — DESIGN ─────────────────────────────────────────────────────────
def s_m7_design(prs):
    slide = blank(prs)
    content_chrome(slide, "MODULE 7 · DESIGN THE PIPELINE", "Apply Phase 1 → Phase 2, Then Lock It")
    card(slide, "PHASE 1 — extract (15 min, whiteboard)", [
        ("Goal: off the mainframe, onto Spring Boot.", False, False, 10, ASC_TEXT),
        ("Inputs: source, copybooks, JCL log, CICS log.", False, False, 10, ASC_TEXT),
        ("Outputs: plan, docs, stories, design, code,", False, False, 10, ASC_TEXT),
        ("           tests, exec report.", False, False, 10, ASC_TEXT),
        ("Constraint: don't rebuild untouched code,", False, False, 10, ASC_TEXT),
        ("           but prove every deletion.", False, False, 10, ASC_TEXT),
        ("Trap (Element 8): you name it after Phase 2.", True, False, 10, ASC_RED),
    ], Inches(0.4), Inches(1.35), Inches(6.1), Inches(3.2), accent=ASC_BLUE)
    card(slide, "PHASE 2 — decompose (25 min)", [
        ("Ask the questions that matter:", True, False, 10.5, ASC_NAVY),
        ("• How do I know what's RUNNING vs what just", False, False, 10, ASC_TEXT),
        ("   EXISTS in the source?", False, False, 10, ASC_TEXT),
        ("• Can a STATIC read of the COBOL be trusted?", False, False, 10, ASC_TEXT),
        ("• What does every agent need to know? (→ KB)", False, False, 10, ASC_TEXT),
        ("• Where does handoff data flow? (→ parent folder)", False, False, 10, ASC_TEXT),
        ("Lock your pipeline before you build.", True, False, 10, ASC_GREEN),
    ], Inches(6.85), Inches(1.35), Inches(6.1), Inches(3.2), accent=ASC_GREEN)
    callout(slide,
            "Instructor will gut-check every team's pipeline before you author anything. The one "
            "question that separates pass from fail: “How will you know what's actually running vs "
            "what just exists in the source?” If you have no step that reads the logs — think again.",
            Inches(0.4), Inches(4.8), Inches(12.55), tone=ASC_AMBER, h=Inches(1.35), label="🚦  GUT-CHECK GATE")
    add_text(slide, "You have a full YEAR of JCL and CICS production logs. They are not decoration.",
             Inches(0.4), Inches(6.35), Inches(12.5), Inches(0.4), size=12, bold=True, italic=True, color=ASC_RED)


def s_m7_challenge(prs):
    slide = blank(prs)
    content_chrome(slide, "MODULE 7 · THE CHALLENGE", "Will You Make the Right Choice?")
    add_text(slide, "Two programs. Both contain dead code. But the dead code hides differently —",
             Inches(0.4), Inches(1.25), Inches(12.5), Inches(0.4), size=13, bold=True, color=ASC_NAVY)
    card(slide, "ACCTPROC — the easy one", [
        ("Its dead paragraphs (5000/6000/7000) are NEVER", False, False, 10.5, ASC_TEXT),
        ("performed by 0000-MAIN-CONTROL.", False, False, 10.5, ASC_TEXT),
        ("", False, False, 4, ASC_TEXT),
        ("→ They're statically UNREACHABLE — orphan", False, False, 10.5, ASC_TEXT),
        ("   paragraphs. A good static parse catches them.", True, False, 10.5, ASC_GREEN),
    ], Inches(0.4), Inches(1.85), Inches(6.1), Inches(2.4), accent=ASC_GREEN)
    card(slide, "CUSTMGMT — the trap 🪤", [
        ("Its dead paragraphs (7000/8000/9000) ARE reachable", False, False, 10.5, ASC_TEXT),
        ("— the EVALUATE dispatches to them on function codes.", False, False, 10.5, ASC_TEXT),
        ("", False, False, 4, ASC_TEXT),
        ("→ Static analysis says KEEP. Only the runtime logs", False, False, 10.5, ASC_TEXT),
        ("   (CARC/CPRG/CMIG = 0 in 2024) prove they're DEAD.", True, False, 10.5, ASC_RED),
    ], Inches(6.85), Inches(1.85), Inches(6.1), Inches(2.4), accent=ASC_RED)
    add_rect(slide, Inches(0.4), Inches(4.5), Inches(12.55), Inches(2.2), fill_rgb=ASC_DARK)
    add_text(slide, "THE LESSON", Inches(0.7), Inches(4.68), Inches(11), Inches(0.35), size=11, bold=True, color=ASC_GREEN)
    add_text(slide,
             "Static and runtime analysis are complementary — and catch DIFFERENT things.\n"
             "A team without a runtime-log-analysis agent WILL migrate CUSTMGMT's dead code.",
             Inches(0.7), Inches(5.05), Inches(11.9), Inches(1.0), size=17, bold=True, color=ASC_WHITE)
    add_text(slide, "That decision — and your evidence for it — is what you're graded on.",
             Inches(0.7), Inches(6.15), Inches(11.9), Inches(0.4), size=12, italic=True, color=ASC_SKY)


# ── MODULE 8 — BUILD ──────────────────────────────────────────────────────────
def s_m8_build(prs):
    slide = blank(prs)
    content_chrome(slide, "MODULE 8 · BUILD VIA MCP", "Author & Push — In Order")
    steps = [
        "REUSE first: find_aava_artifact the AzureBlob tools + Current DateTime Tool. Don't recreate.",
        "KB: create_aava_knowledge_base from cobol-modernization-standards.txt (or reuse if it exists).",
        "Agents: create_aava_agent for each — STEP / DO-NOT-PROCEED instructions. Suffix names w/ your initials.",
        "Workflow: create_aava_workflow — remember workflowConfig: {}  and the chaining prompt on agents 2+.",
        "Thread azure_parent_folder from Agent 1 through every downstream agent (Commandment #2).",
    ]
    for i, s in enumerate(steps):
        step_box(slide, i + 1, s, Inches(0.4), Inches(1.45) + i * Inches(0.7), w=Inches(12.55), tone=ASC_GREEN)
    add_rect(slide, Inches(0.4), Inches(5.1), Inches(12.55), Inches(1.6), fill_rgb=ASC_LIGHT, line_rgb=ASC_BORDER, line_width_pt=0.5)
    add_text(slide, "CREATION ORDER (Commandment #10)", Inches(0.7), Inches(5.28), Inches(12), Inches(0.35), size=11, bold=True, color=ASC_DARK)
    add_text(slide, "KB  →  GR  →  Tool  →  Agent  →  Workflow",
             Inches(0.7), Inches(5.68), Inches(12), Inches(0.6), size=22, bold=True, color=ASC_NAVY)
    add_text(slide, "Build bottom-up — each layer references the one before it.",
             Inches(0.7), Inches(6.3), Inches(12), Inches(0.35), size=11, italic=True, color=ASC_MUTED)


def s_m8_reference(prs):
    slide = blank(prs)
    content_chrome(slide, "MODULE 8 · REFERENCE SHAPE", "The Production 8-Agent Pipeline")
    agents = [
        ("1", "COBOL Source Parser", "static structure + dependency map", ASC_BLUE),
        ("2", "Runtime Log Analyzer ⭐", "ACTIVE vs DEAD, with evidence", ASC_RED),
        ("3", "Tech Doc Creator", "document ACTIVE code only", ASC_BLUE),
        ("4", "User Story Creator", "stories for ACTIVE functionality", ASC_BLUE),
        ("5", "HLD / LLD Creator", "Spring Boot target design", ASC_GREEN),
        ("6", "Spring Boot Code Gen", "Java for ACTIVE logic only", ASC_GREEN),
        ("7", "JUnit Test Generator", "JUnit 5 + Mockito coverage", ASC_GREEN),
        ("8", "Executive Report Gen", "board-ready HTML + savings %", ASC_NAVY),
    ]
    for i, (n, name, desc, col) in enumerate(agents):
        col_i = i % 2
        row = i // 2
        x = Inches(0.4) + col_i * Inches(6.45)
        y = Inches(1.4) + row * Inches(1.05)
        add_rect(slide, x, y, Inches(6.2), Inches(0.92), fill_rgb=ASC_GRAY, line_rgb=ASC_BORDER, line_width_pt=0.5)
        add_rect(slide, x, y, Inches(0.6), Inches(0.92), fill_rgb=col)
        add_text(slide, n, x, y + Inches(0.26), Inches(0.6), Inches(0.4), size=18, bold=True, color=ASC_WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, name, x + Inches(0.75), y + Inches(0.13), Inches(5.3), Inches(0.35), size=12, bold=True, color=ASC_DARK)
        add_text(slide, desc, x + Inches(0.75), y + Inches(0.5), Inches(5.3), Inches(0.32), size=9.5, italic=True, color=ASC_NAVY)
    callout(slide,
            "Agent 2 is the whole point. Without a dedicated runtime-evidence step you can't defensibly "
            "retire code — and every downstream agent is scoped to ACTIVE CODE ONLY. That scoping is what "
            "turns Agent 2's intelligence into real savings. Your shape can be 6–9 agents — but Parse + "
            "Runtime-Analysis + active-scoped downstream + Exec-Report must all be there.",
            Inches(0.4), Inches(5.7), Inches(12.55), tone=ASC_RED, h=Inches(1.15), label="⭐  WHY AGENT 2 MATTERS")


# ── MODULE 9 — TEST ───────────────────────────────────────────────────────────
def s_m9_test(prs):
    slide = blank(prs)
    content_chrome(slide, "MODULE 9 · TEST, DEBUG, VALIDATE", "Unit-Test Before You Trust")
    steps = [
        "execute_single_agent on Agent 1 → confirm it parsed structure + generated the parent folder.",
        "execute_single_agent on Agent 2 → THE validation moment (see below).",
        "trigger_workflow → poll_workflow_result / stream_workflow_progress_formatted.",
        "Debug: synthetic data? → Commandment #7. Tighten with blocking language, re-run.",
    ]
    for i, s in enumerate(steps):
        step_box(slide, i + 1, s, Inches(0.4), Inches(1.4) + i * Inches(0.66), w=Inches(12.55), tone=ASC_AMBER)
    add_rect(slide, Inches(0.4), Inches(4.2), Inches(12.55), Inches(2.5), fill_rgb=ASC_DARK)
    add_text(slide, "THE VALIDATION MOMENT — did Agent 2 actually read the logs?", Inches(0.7), Inches(4.4), Inches(12), Inches(0.4),
             size=12, bold=True, color=ASC_GREEN)
    add_lines(slide, [
        ("✅  PASS: output is a dead-code inventory — each dead paragraph with its exec count (0),", False, False, 11.5, ASC_WHITE),
        ("        last-executed date, and a REASON (replaced by DB2, blocked by compliance, etc.).", False, False, 11.5, ASC_WHITE),
        ("", False, False, 4, ASC_TEXT),
        ("❌  FAIL: vague summary with no counts / dates / reasons = it never read the logs.", False, False, 11.5, ASC_SKY),
        ("        Add 'YOU MUST read both logs... DO NOT PROCEED until...' and re-run.", False, False, 11.5, ASC_SKY),
    ], Inches(0.7), Inches(4.85), Inches(11.9), Inches(1.7), spacing=18)


# ── MODULE 10 — APPROVE + DEBRIEF ─────────────────────────────────────────────
def s_m10_approve(prs):
    slide = blank(prs)
    content_chrome(slide, "MODULE 10 · SHIP IT", "Approve One-by-One, Then Defend Your Calls")
    card(slide, "APPROVE — in order, after validation", [
        ("approve_aava_artifact for each artifact that", False, False, 10.5, ASC_TEXT),
        ("passed its test — KB → tools → agents.", False, False, 10.5, ASC_TEXT),
        ("", False, False, 4, ASC_TEXT),
        ("Never approve what you haven't seen produce", False, False, 10.5, ASC_TEXT),
        ("correct output. Body uses  id  not agentId.", True, False, 10.5, ASC_NAVY),
    ], Inches(0.4), Inches(1.35), Inches(6.1), Inches(2.4), accent=ASC_GREEN)
    card(slide, "DEBRIEF — the headline question", [
        ("“How many paragraphs did you RETIRE,", True, False, 12, ASC_DARK),
        ("and what's your EVIDENCE?”", True, False, 12, ASC_DARK),
        ("", False, False, 4, ASC_TEXT),
        ("Be ready to defend every deletion with the", False, False, 10.5, ASC_TEXT),
        ("0-exec counts + stale dates from the logs.", False, False, 10.5, ASC_TEXT),
    ], Inches(6.85), Inches(1.35), Inches(6.1), Inches(2.4), accent=ASC_BLUE)
    callout(slide,
            "Instructor reveals the answer key only AFTER each team presents its numbers. Let teams "
            "commit to a migrate/retire count first — the discussion is where the learning lands.",
            Inches(0.4), Inches(4.05), Inches(12.55), tone=ASC_AMBER, h=Inches(0.95), label="🚦  INSTRUCTOR NOTE")
    add_text(slide, "Next slide: the answer key (instructor reveal).",
             Inches(0.4), Inches(5.25), Inches(12), Inches(0.4), size=12, italic=True, color=ASC_MUTED)


def s_m10_answer(prs):
    slide = blank(prs)
    content_chrome(slide, "MODULE 10 · ANSWER KEY  (INSTRUCTOR REVEAL)", "9 Dead · 23 Active · 28.9% Saved")
    add_text(slide, "Reveal only after teams have presented their own numbers.",
             Inches(0.4), Inches(1.2), Inches(12.5), Inches(0.35), size=10.5, bold=True, italic=True, color=ASC_RED)
    # CUSTMGMT
    card(slide, "CUSTMGMT — 12 active / 4 DEAD", [
        ("7000 ARCHIVE-CUSTOMERS  ·  last 2019  ·  DB2 took over", False, False, 9.5, ASC_TEXT),
        ("7100 WRITE-ARCHIVE      ·  child of 7000", False, False, 9.5, ASC_TEXT),
        ("8000 PURGE-INACTIVE     ·  2018  ·  compliance blocked", False, False, 9.5, ASC_TEXT),
        ("9000 MIGRATE-TO-NEW-SYS ·  NEVER  ·  cancelled 2017", False, False, 9.5, ASC_TEXT),
        ("🪤 reachable in code — only logs prove dead", True, False, 9.5, ASC_RED),
    ], Inches(0.4), Inches(1.65), Inches(6.1), Inches(2.3), accent=ASC_RED)
    card(slide, "ACCTPROC — 11 active / 5 DEAD", [
        ("5000 CALC-MONTHLY-FEES  ·  2021  ·  FEECALC took over", False, False, 9.5, ASC_TEXT),
        ("5100 LOW-BALANCE-FEE    ·  5200 EXCESS-TRAN-FEE", False, False, 9.5, ASC_TEXT),
        ("6000 PAPER-STATEMENTS   ·  2020  ·  e-statements", False, False, 9.5, ASC_TEXT),
        ("7000 Y2K-DATE-CHECK     ·  2000  ·  obsolete", False, False, 9.5, ASC_TEXT),
        ("unreachable — static parse catches these", True, False, 9.5, ASC_GREEN),
    ], Inches(6.85), Inches(1.65), Inches(6.1), Inches(2.3), accent=ASC_GREEN)
    # scorecard
    add_rect(slide, Inches(0.4), Inches(4.15), Inches(12.55), Inches(2.55), fill_rgb=ASC_DARK)
    add_text(slide, "COMBINED — what the executive report must show", Inches(0.7), Inches(4.32), Inches(12), Inches(0.4),
             size=12, bold=True, color=ASC_GREEN)
    stats = [
        ("32", "total paragraphs"),
        ("23", "ACTIVE (71.9%)"),
        ("9", "DEAD (28.1%)"),
        ("~140", "dead lines (28.9%)"),
        ("28.9%", "scope reduction"),
    ]
    for i, (big, lbl) in enumerate(stats):
        x = Inches(0.7) + i * Inches(2.5)
        add_rect(slide, x, Inches(4.85), Inches(2.3), Inches(1.5), fill_rgb=RGBColor(0x18, 0x36, 0x58), line_rgb=ASC_BLUE, line_width_pt=0.75)
        add_text(slide, big, x, Inches(5.05), Inches(2.3), Inches(0.7), size=30, bold=True, color=ASC_WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, lbl, x, Inches(5.75), Inches(2.3), Inches(0.5), size=10, color=ASC_SKY, align=PP_ALIGN.CENTER)
    add_text(slide, "vs the previous vendor's 1-for-1 rewrite — ~29% of scope cut, with evidence the board can defend.",
             Inches(0.4), Inches(6.78), Inches(12.5), Inches(0.3), size=9.5, italic=True, color=ASC_MUTED)


# ── CLOSE ─────────────────────────────────────────────────────────────────────
def s_close(prs):
    slide = blank(prs)
    set_bg(slide, ASC_DARK)
    add_rect(slide, Inches(0), Inches(0), Inches(0.18), SLIDE_H, fill_rgb=ASC_GREEN)
    add_rect(slide, Inches(0.18), Inches(0), SLIDE_W - Inches(0.18), Inches(0.06), fill_rgb=ASC_BLUE)
    add_text(slide, "You Built It.", Inches(0.5), Inches(0.9), Inches(12), Inches(1.0), size=48, bold=True, color=ASC_WHITE)
    add_text(slide,
             "You took a raw client problem and shipped a governed, tested, multi-agent workflow —\n"
             "and you retired dead code with evidence the board can defend.",
             Inches(0.5), Inches(2.0), Inches(12.3), Inches(1.0), size=16, color=ASC_SKY)
    achievements = [
        ("Decomposed", "one agent per critical task"),
        ("Authored via MCP", "KB → tools → agents → workflow"),
        ("Validated", "tested each artifact before trust"),
        ("Approved 1-by-1", "governance discipline"),
        ("Proved the savings", "28.9% scope cut, evidence-backed"),
    ]
    for i, (a, d) in enumerate(achievements):
        y = Inches(3.3) + i * Inches(0.66)
        add_rect(slide, Inches(0.5), y, Inches(8.0), Inches(0.56), fill_rgb=RGBColor(0x18, 0x36, 0x58))
        add_rect(slide, Inches(0.5), y, Inches(0.07), Inches(0.56), fill_rgb=ASC_GREEN)
        add_text(slide, "✓  " + a, Inches(0.7), y + Inches(0.13), Inches(3.2), Inches(0.32), size=12, bold=True, color=ASC_WHITE)
        add_text(slide, d, Inches(3.9), y + Inches(0.14), Inches(4.4), Inches(0.32), size=10.5, color=ASC_SKY)
    add_rect(slide, Inches(8.9), Inches(3.3), Inches(4.0), Inches(3.2), fill_rgb=RGBColor(0x15, 0x32, 0x52), line_rgb=ASC_BLUE, line_width_pt=0.75)
    add_text(slide, "THE ASCENDION EDGE", Inches(9.1), Inches(3.55), Inches(3.6), Inches(0.4), size=10, bold=True, color=ASC_GREEN, align=PP_ALIGN.CENTER)
    add_text(slide,
             "Runtime intelligence + AI-native\nengineering turns a $2.4M\n1-for-1 rewrite into a\ndefensible, scoped\nmodernization.\n\nThat's what clients buy.",
             Inches(9.1), Inches(4.0), Inches(3.6), Inches(2.3), size=12, color=ASC_WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, FOOTER, Inches(0.5), Inches(7.12), Inches(10), Inches(0.3), size=7, color=RGBColor(0x50, 0x70, 0x90))


def s_reference(prs):
    slide = blank(prs)
    content_chrome(slide, "APPENDIX", "Quick-Reference & Resources")
    card(slide, "MCP CHEAT-SHEET", [
        ("Verify:  test_aava_connection", False, False, 9.5, ASC_TEXT),
        ("Exists?  find_aava_artifact / list_aava_artifacts", False, False, 9.5, ASC_TEXT),
        ("Create:  create_aava_{kb,guardrail,tool,agent,workflow}", False, False, 9.5, ASC_TEXT),
        ("Test:    execute_single_agent / execute_single_tool", False, False, 9.5, ASC_TEXT),
        ("Run:     trigger_workflow → poll_workflow_result", False, False, 9.5, ASC_TEXT),
        ("Approve: approve_aava_artifact  (one at a time)", False, False, 9.5, ASC_TEXT),
        ("Modify:  clone_aava_artifact  (APPROVED is immutable)", False, False, 9.5, ASC_TEXT),
    ], Inches(0.4), Inches(1.35), Inches(7.0), Inches(3.2), accent=ASC_BLUE)
    card(slide, "REMEMBER", [
        ("KB→GR→Tool→Agent→Workflow", True, False, 11, ASC_NAVY),
        ("Thread azure_parent_folder", False, False, 10, ASC_TEXT),
        ("Guardrails on agents, not workflows", False, False, 10, ASC_TEXT),
        ("workflowConfig: {} on create", False, False, 10, ASC_TEXT),
        ("DO NOT PROCEED on tool calls", False, False, 10, ASC_TEXT),
        ("Reuse > create", True, False, 10, ASC_GREEN),
    ], Inches(7.6), Inches(1.35), Inches(5.35), Inches(3.2), accent=ASC_GREEN)
    add_rect(slide, Inches(0.4), Inches(4.75), Inches(12.55), Inches(1.95), fill_rgb=ASC_LIGHT, line_rgb=ASC_BORDER, line_width_pt=0.5)
    add_text(slide, "RESOURCES", Inches(0.7), Inches(4.9), Inches(12), Inches(0.35), size=11, bold=True, color=ASC_DARK)
    add_lines(slide, [
        ("Master playbook:  AAVA-Mastery / AAVAAdvancedProblemSolvingSystem.md  (11 Commandments, 8 phases)", False, False, 10.5, ASC_TEXT),
        ("MCP server:       github.com/wgpullen/AAVA-MCP-Server  (63 tools)", False, False, 10.5, ASC_TEXT),
        ("This lab:         github.com/wgpullen/AAVA-Advanced-Training  (capstone + solution key + this deck)", False, False, 10.5, ASC_TEXT),
        ("Platform:         https://int-ai.aava.ai   ·   Realm 32   ·   Support: ascendionava.support@ascendion.com", False, False, 10.5, ASC_TEXT),
    ], Inches(0.7), Inches(5.3), Inches(12.2), Inches(1.3), spacing=18)


# ══════════════════════════════════════════════════════════════════════════════
def build(output_path):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    s_cover(prs)
    s_agenda(prs)
    s_audience(prs)

    section_header(prs, "MODULE 0", "Environment Bring-up", "Get everyone green before we build")
    s_m0_setup(prs)
    s_m0_surfaces(prs)

    section_header(prs, "MODULE 1", "Mental Model + 11 Commandments", "Understand the platform before you build")
    s_m1_glance(prs)
    s_m1_commandments_a(prs)
    s_m1_commandments_b(prs)

    section_header(prs, "MODULE 2", "Problem Decomposition", "One agent = one independent task")
    s_m2_golden(prs)
    s_m2_phase1(prs)
    s_m2_phase2(prs)

    section_header(prs, "MODULE 3", "KBs, Guardrails, Tools", "The building blocks agents stand on")
    s_m3_kb_gr(prs)
    s_m3_tools(prs)

    section_header(prs, "MODULE 4", "Agent Creation via Repo + MCP", "Author and push artifacts programmatically")
    s_m4_template(prs)
    s_m4_configs_mcp(prs)

    section_header(prs, "MODULE 5", "Governance Discipline", "Test → Validate → Approve one-by-one")
    s_m5_governance(prs)

    s_lunch(prs)

    section_header(prs, "MODULE 6", "Capstone Briefing", "Meridian National Bank — a real client problem")
    s_m6_client(prs)
    s_m6_artifacts(prs)

    section_header(prs, "MODULE 7", "Design the Pipeline", "Apply Phase 1–2, then lock it")
    s_m7_design(prs)
    s_m7_challenge(prs)

    section_header(prs, "MODULE 8", "Build via MCP", "KB → GR → Tool → Agent → Workflow")
    s_m8_build(prs)
    s_m8_reference(prs)

    section_header(prs, "MODULE 9", "Test, Debug, Validate", "Unit-test before you trust")
    s_m9_test(prs)

    section_header(prs, "MODULE 10", "Approve & Debrief", "Ship it, then defend your dead-code calls")
    s_m10_approve(prs)
    s_m10_answer(prs)

    s_close(prs)
    s_reference(prs)

    prs.save(output_path)
    print(f"✅ Saved {output_path} — {len(prs.slides.__iter__.__self__._sldIdLst)} slides")


if __name__ == "__main__":
    import sys
    out = sys.argv[1] if len(sys.argv) > 1 else "AAVA_Advanced_Engineering_Lab_v1.pptx"
    build(out)
